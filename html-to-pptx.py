"""presentation.html → presentation.pptx (16:9, 다크 톤, 자동 매핑)

각 슬라이드:
  - 상단 메타 (회색 작은 글씨)
  - 큰 제목 (h1)
  - 본문 (lede + bullet list)
  - 표 (있으면 native pptx 표로)
  - 캡처 이미지 (백업 A5-A8)
  - 발표자 노트 (narration.json 텍스트)
"""
import json
import os
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent
HTML_PATH = ROOT / 'presentation.html'

# ─── color tokens ───
BG     = RGBColor(0x0c, 0x0b, 0x0a)
FG     = RGBColor(0xf6, 0xf2, 0xea)
DIM    = RGBColor(0xaa, 0xa3, 0x9a)
MUTE   = RGBColor(0x6f, 0x6a, 0x62)
ACCENT = RGBColor(0xff, 0x48, 0x36)
ACCENT_SOFT = RGBColor(0xf4, 0xc8, 0x69)
LINE   = RGBColor(0x2a, 0x27, 0x24)


def set_slide_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_text(slide, x_in, y_in, w_in, h_in, text, *,
             size=14, bold=False, color=FG, align=None, italic=False):
    """텍스트 박스 추가."""
    tx = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    # paragraphs
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if align == 'center': p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = 'Pretendard'
    return tx


def add_bullets(slide, x_in, y_in, w_in, h_in, items, *,
                size=13, color=DIM, bullet='• '):
    tx = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet + txt
        p.space_after = Pt(4)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = 'Pretendard'
    return tx


def add_native_table(slide, x_in, y_in, w_in, h_in, rows_data, *,
                     header_color=ACCENT, body_color=DIM, header_bg=None):
    """rows_data: list of list. first row = header."""
    if not rows_data: return None
    n_rows = len(rows_data)
    n_cols = max(len(r) for r in rows_data)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tbl = table_shape.table
    for ri, row in enumerate(rows_data):
        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            txt = row[ci] if ci < len(row) else ''
            cell.text = str(txt)
            # 배경
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG if ri > 0 else RGBColor(0x15, 0x13, 0x0f)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11 if ri > 0 else 10)
                    run.font.bold = (ri == 0)
                    run.font.color.rgb = header_color if ri == 0 else body_color
                    run.font.name = 'Pretendard'
    return table_shape


def add_card(slide, x_in, y_in, w_in, h_in, title, body, *, bordered=False):
    """카드형 박스 (배경 + 테두리)."""
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0x15, 0x13, 0x0f)
    shp.line.color.rgb = ACCENT if bordered else LINE
    shp.line.width = Pt(1.5 if bordered else 0.5)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = tf.margin_bottom = Inches(0.12)
    p1 = tf.paragraphs[0]
    p1.text = title
    for run in p1.runs:
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        run.font.name = 'Pretendard'
    p2 = tf.add_paragraph()
    p2.text = body
    p2.space_before = Pt(4)
    for run in p2.runs:
        run.font.size = Pt(11)
        run.font.color.rgb = DIM
        run.font.name = 'Pretendard'
    return shp


def extract_clean_text(el):
    """공백·줄바꿈 정리해서 텍스트 추출."""
    if not el: return ''
    return ' '.join(el.get_text(' ', strip=True).split())


def main():
    soup = BeautifulSoup(HTML_PATH.read_text(encoding='utf-8'), 'html.parser')

    narrations = {}
    nj = ROOT / 'narration.json'
    if nj.exists():
        for s in json.loads(nj.read_text(encoding='utf-8')):
            narrations[s['index']] = s['text']

    sections = soup.find_all('section', class_='slide')
    print(f'▶ {len(sections)} slides')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, sec in enumerate(sections, 1):
        is_backup = 'backup' in sec.get('class', [])
        label = f'A{idx-10}' if is_backup else str(idx)
        print(f'  → slide {label}')

        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide, BG)

        # 메타 (상단)
        meta = sec.find(class_='meta')
        if meta:
            mtxt = extract_clean_text(meta)
            add_text(slide, 0.5, 0.25, 12, 0.35, mtxt, size=10, color=MUTE)

        # h1
        h1 = sec.find('h1')
        if h1:
            add_text(slide, 0.5, 0.6, 12, 1.0, extract_clean_text(h1),
                     size=32, bold=True, color=FG)

        # lede
        lede = sec.find(class_='lede')
        lede_y = 1.5
        if lede:
            add_text(slide, 0.5, lede_y, 12, 0.7, extract_clean_text(lede),
                     size=15, color=DIM, italic=True)
            lede_y += 0.8

        body_y = lede_y + 0.2

        # 캡처 이미지 (A5~A8 백업) — 좌측에
        imgs = sec.find_all('img', class_='shot')
        body_x = 0.5
        body_w = 12.3
        if imgs:
            for i, img in enumerate(imgs[:2]):
                src = img.get('src', '')
                if src.startswith('./'): src = src[2:]
                img_path = ROOT / src
                if img_path.exists():
                    img_y = body_y + i * 2.7
                    try:
                        slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(img_y),
                                                  width=Inches(5.5))
                    except Exception as e:
                        print(f'    ! img fail: {e}')
            body_x = 6.2
            body_w = 6.6

        # 표
        tables = sec.find_all('table')
        cards = sec.find_all(class_='card')
        flow_nodes = sec.find_all(class_='node')
        # quote
        quotes = sec.find_all(class_='quote')

        # 우선순위 — table → card grid → flow → quote → bullet list
        if tables:
            for ti, tbl in enumerate(tables[:1]):
                rows = []
                for tr in tbl.find_all('tr'):
                    cells = [extract_clean_text(td) for td in tr.find_all(['th', 'td'])]
                    rows.append(cells)
                if rows:
                    tbl_h = min(5.2, 0.4 * len(rows))
                    add_native_table(slide, body_x, body_y, body_w, tbl_h, rows)
                    body_y += tbl_h + 0.2
        elif cards:
            n = min(len(cards), 4)
            card_w = (body_w - (n - 1) * 0.2) / n
            for ci, c in enumerate(cards[:n]):
                title = ''
                body_lines = []
                for el in c.find_all(['h2', 'h3'], limit=1):
                    title = extract_clean_text(el)
                for el in c.find_all(['p', 'li']):
                    t = extract_clean_text(el)
                    if t and len(t) < 200: body_lines.append(t)
                add_card(slide, body_x + ci * (card_w + 0.2), body_y, card_w, 4,
                         title or f'카드 {ci+1}', '\n'.join(body_lines[:3]),
                         bordered=('bordered' in c.get('class', [])))
            body_y += 4.2
        elif flow_nodes:
            n = len(flow_nodes)
            node_w = (body_w - (n - 1) * 0.15) / n
            for ni, node in enumerate(flow_nodes[:n]):
                title = ''
                sub = ''
                strong = node.find(class_='strong')
                subel = node.find(class_='sub')
                title = extract_clean_text(strong) if strong else ''
                sub = extract_clean_text(subel) if subel else ''
                is_ml = 'ml' in node.get('class', [])
                add_card(slide, body_x + ni * (node_w + 0.15), body_y, node_w, 2,
                         title, sub, bordered=is_ml)
            body_y += 2.2

        # 본문 bullets (lede + table/card 이후 남은 텍스트)
        if not (tables or cards or flow_nodes):
            items = []
            for el in sec.find_all(['li'], recursive=True):
                if el.find_parent(class_='card') or el.find_parent(class_='node'): continue
                t = extract_clean_text(el)
                if t and t not in items and 5 < len(t) < 250:
                    items.append(t)
            if items:
                add_bullets(slide, body_x, body_y, body_w, 4.5, items[:12], size=12)

        # quote (하단)
        if quotes:
            qtxt = extract_clean_text(quotes[0])
            if qtxt:
                add_text(slide, 0.5, 6.5, 12.3, 0.8,
                         '" ' + qtxt[:300] + ('...' if len(qtxt) > 300 else '') + ' "',
                         size=11, color=ACCENT_SOFT, italic=True)

        # 발표자 노트 (narration)
        nar = narrations.get(idx, '')
        if nar:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = nar

    out = ROOT / 'presentation.pptx'
    prs.save(out)
    print(f'\n✓ saved {out.name}  ({out.stat().st_size//1024}KB)')


if __name__ == '__main__':
    main()
